import pandas as pd
import glob
import os
import matplotlib.pyplot as plt

def detect_outliers_iqr(df, column):
    """Detect outliers using the interquartile range method."""
    Q1 = df[column].quantile(0.25)
    Q3 = df[column].quantile(0.75)
    IQR = Q3 - Q1
    lower_bound = Q1 - 1.5 * IQR
    upper_bound = Q3 + 1.5 * IQR
    outliers = df[(df[column] < lower_bound) | (df[column] > upper_bound)]
    df_clean = df[(df[column] > lower_bound) & (df[column] < upper_bound)]
    return outliers, df_clean

def clean_sauts(df, column):
    """Nettoyage des sauts d'index. La variable doit etre un index
    Utilise la methode de la derivée
    Contrainte :
    -La donnée doit etre en kWh ou en m3
    -Minimum une journée de data au 1/4 heure"""

    # On calcul la dérivée de l'index, ce qui correspond a une puissance ou un debit
    delta_time = df.index.to_series().diff().dt.total_seconds()
    df['derivative'] = df['tagValue'].diff()*3600 / delta_time

    df['rolling_min'] = df['tagValue'].rolling(5).min()
    df['ecart_rel_rol_min'] = abs(100*(df['tagValue']-df['rolling_min'])/df['rolling_min'])

    # On supprime les valeurs superieures a 4 000 kW (4 MW)
    #  ou 4000 (m3*3600)/s
    df_clean = df[(abs(df['derivative']) < 4000) & (df['ecart_rel_rol_min'] < 50) ].copy()
    outliers = df[(abs(df['derivative']) > 4000) | (df['ecart_rel_rol_min'] > 50) ].copy()

    # On enleve les colonnes supplementaires
    df_clean.drop(columns=['derivative','rolling_min','ecart_rel_rol_min'], inplace=True)
    outliers.drop(columns=['derivative', 'rolling_min', 'ecart_rel_rol_min'], inplace=True)

    return outliers, df_clean

def clean_zeroes(df, column):

    df_zeroes = df.loc[df[column] == 0]
    df_out = df.loc[df[column] != 0]
    return df_zeroes, df_out

def clean_stuck_index(df, column):

    df_tmp = df['tagValue'].diff().rolling(5).mean()
    df_clean = df_tmp.loc[df_tmp[column] != 0]
    df_blok = df_tmp.loc[df_tmp[column] == 0]
    return df_blok, df_clean

def create_pdf(file_name, figure):
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.pdfgen import canvas

    projet_name = 'Projet_test'

    # Create a canvas with landscape orientation
    c = canvas.Canvas(file_name, pagesize=landscape(letter))

    # Add text
    from datetime import datetime
    current_time = datetime.now()
    current_time= current_time.strftime("%d/%m/%Y %H:%M")

    c.setFont('Times-Bold',28)

    # Page de titre
    c.drawImage('C:\\NETTOYEUR\\img\\EM_Logo.png',260,300,width=1.28*170,height=170,mask='auto')
    c.drawString(260, 250, 'Rapport de nettoyage des données')
    c.drawString(260, 200, current_time)
    c.drawString(260, 150, projet_name)

    # Nouvelle page
    c.showPage()

    # Add image (figure)
    img_path = figure  # Replace with your image file path
    c.drawImage(img_path, 100, 100, width=500, height=400)

    # Close the canvas
    c.save()

    print('pdf cree')

def is_index(df):
    """Test if dataset is index
    Input : dataframe with time in 1st column and value in 2nd"""
    ## TODO : rentre test plus robuste : avec remise a zero compteurs
    # On commence par ecarter les valeurs pourries
    # On vire eventuellement une colonne tagName
    if 'quality' in df.columns:
        df = df.drop(['quality'], axis=1, errors='ignore')
    if 'tagName' in df.columns:
        df = df.drop(['tagName'], axis=1, errors='ignore')

    outliers, df = detect_outliers_iqr(df,df.columns[0])
    if not pd.api.types.is_datetime64_any_dtype(df.index):
        df = df.set_index(pd.to_datetime(df[df.columns[0]], format='%Y-%m-%d %H:%M:%S'))
    if df.columns[0] == 'ts':
        df = df.drop([df.columns[0]], axis=1, errors='ignore')

    if df[df.columns[0]].is_monotonic_increasing and df.empty == False:
        return True
    else:
        return False

def consumption_from_index(df):
    """Compute consumption from index and return consumption
    by day, month, year"""

    df_diff = df.diff()

    df_by_days = df_diff.resample('D').sum()
    df_by_days.index = df_by_days.index.strftime('%Y-%m-%d')

    df_by_months = df_diff.resample('m').sum()
    df_by_months.index = df_by_months.index.strftime('%B-%Y')

    df_by_years = df_diff.resample('Y').sum()
    df_by_years.index = df_by_years.index.strftime('%Y')

    conso = {'cons_brut':df_diff,
            'cons_day':df_by_days,
             'cons_month':df_by_months,
             'cons_year':df_by_years}

    return conso

def img_consumption(conso, dir):
    '''Utilise la sortie de la fonction consumption_from_index'''

    # Courbe de la consommation brute
    df = conso['cons_brut']
    fig, ax1, = plt.subplots(figsize=(8, 6))
    ax1.bar(df.index, df[df.columns[0]])
    ax1.xaxis.grid(True, linestyle='--')
    ax1.legend(loc='best')
    ax1.set_title('Consommation brute')
    ax1.set_ylabel('Consommation (kWh)')
    plt.tick_params(axis='x', labelsize='medium')
    plt.savefig(dir + 'tmp_brut.png')

    # Courbe de la consommation Jours
    df = conso['cons_day']
    fig, ax1, = plt.subplots(figsize=(8, 6))
    ax1.bar(df.index, df[df.columns[0]])
    ax1.xaxis.grid(True, linestyle='--')
    ax1.legend(loc='best')
    ax1.set_title('Consommation par jour')
    ax1.set_ylabel('Consommation (kWh)')
    plt.tick_params(axis='x', labelsize='medium')
    plt.savefig(dir + 'tmp_day.png')
    plt.close(fig)

    # Courbe de la consommation Mois
    df = conso['cons_month']
    fig, ax1, = plt.subplots(figsize=(8, 6))
    ax1.bar(df.index, df[df.columns[0]])
    ax1.xaxis.grid(True, linestyle='--')
    ax1.legend(loc='best')
    ax1.set_title('Consommation par mois')
    ax1.set_ylabel('Consommation (kWh)')
    plt.tick_params(axis='x', labelsize='medium')
    plt.savefig(dir + 'tmp_month.png')
    plt.close(fig)

    # Courbe de la consommation Annee
    df = conso['cons_year']
    fig, ax1, = plt.subplots(figsize=(8, 6))
    ax1.bar(df.index, df[df.columns[0]])
    ax1.xaxis.grid(True, linestyle='--')
    ax1.legend(loc='best')
    ax1.set_title('Consommation par année')
    ax1.set_ylabel('Consommation (kWh)')
    plt.tick_params(axis='x', labelsize='medium')
    plt.savefig(dir + 'tmp_year.png')
    plt.close(fig)

def xl_consumption(liste_conso, dir):
    '''Utilise la sortie de la fonction consumption_from_index'''
    from openpyxl import load_workbook
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    # INITIALISATION POWERPOINT
    prs = Presentation()

    for conso in liste_conso:
        var_name = conso['cons_brut'].columns[0]

        # On check d'abord le range des dates
        total_delta = conso['cons_brut'].index[-1] - conso['cons_brut'].index[0]

        # On initialise les tests
        ae_brut, ae_days, ae_month, ae_year = False,False, False, False

        # Si l'interval est inferieure a un mois, on autorise les graphes data brut et par jour
        if total_delta.days <= 31:
            ae_brut = True
            ae_days = True
        else:
            ae_month = True
            ae_year = True

        # SI JAMAIS LES DATA SONT TRES PEU DENSES (moins de 1 valeur/jour) MAIS SUR UN LONG INTERVAL DE TEMPS
        # On autorise les analyses longueurs
        if total_delta.days > 31 and total_delta.days > len(conso['cons_brut'].index):
            ae_month = True
            ae_year = True

        # TRACE DES DONNES BRUTES EN TREND DANS TOUS LES CAS
        if ae_brut:

            # ON commence par retirer les NaN
            conso['cons_brut'] = conso['cons_brut'].dropna(how='any')

            # Add a slide
            slide_layout = prs.slide_layouts[5]  # 5 is the index for Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            slide.placeholders[0].text = var_name.replace('value', '').replace('Value', '')

            # Define chart data
            chart_data = CategoryChartData()
            chart_data.categories = conso['cons_brut'].index.to_list()
            chart_data.add_series("Consommation Complète (kWh)", conso['cons_brut'][conso['cons_brut'].columns[0]])

            # Define chart type and position
            x, y, cx, cy = Inches(0.5), Inches(2), Inches(9), Inches(5)

            # Add a chart to the slide
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
            ).chart

            # Set chart title
            chart.has_title = True
            chart.chart_title.text_frame.text = "Consommation brute (kWh)"

            # Set axes titles
            category_axis = chart.category_axis
            category_axis.has_major_gridlines = False
            category_axis.tick_labels.font.bold = False

            value_axis = chart.value_axis
            value_axis.has_major_gridlines = False
            value_axis.tick_labels.font.bold = False
            value_axis.tick_labels.size = 12

        if ae_days:
            # Add a slide
            slide_layout = prs.slide_layouts[5]  # 5 is the index for Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            slide.placeholders[0].text = var_name.replace('value', '').replace('Value', '')

            # Define chart data
            chart_data = CategoryChartData()
            chart_data.categories = conso['cons_day'].index.to_list()
            chart_data.add_series("Consommation Journalière (kWh)", conso['cons_day'][conso['cons_day'].columns[0]])

            # Define chart type and position
            x, y, cx, cy = Inches(0.5), Inches(2), Inches(9), Inches(5)

            # Add a chart to the slide
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            # Set chart title
            chart.has_title = True
            chart.chart_title.text_frame.text = "Consommation Journalière (kWh)"

            # Set axes titles
            category_axis = chart.category_axis
            category_axis.has_major_gridlines = True
            category_axis.tick_labels.font.bold = True

            value_axis = chart.value_axis
            value_axis.has_major_gridlines = True
            value_axis.tick_labels.font.bold = True
            value_axis.tick_labels.size = 12

        if ae_month:
            # Add a slide
            slide_layout = prs.slide_layouts[5]  # 5 is the index for Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            slide.placeholders[0].text = var_name.replace('value','').replace('Value','')

            # Define chart data
            chart_data = CategoryChartData()
            chart_data.categories = conso['cons_month'].index.to_list()
            chart_data.add_series("Consommation Mensuelle (kWh)", conso['cons_month'][conso['cons_month'].columns[0]])

            # Define chart type and position
            x, y, cx, cy = Inches(0.5), Inches(2), Inches(9), Inches(5)

            # Add a chart to the slide
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            # Set chart title
            chart.has_title = True
            chart.chart_title.text_frame.text = "Consommation Mensuelle (kWh)"

            # Set axes titles
            category_axis = chart.category_axis
            category_axis.has_major_gridlines = True
            category_axis.tick_labels.font.bold = False

            value_axis = chart.value_axis
            value_axis.has_major_gridlines = True
            value_axis.tick_labels.font.bold = False
            value_axis.tick_labels.size = 12

        if ae_year:
            # Add a slide
            slide_layout = prs.slide_layouts[5]  # 5 is the index for Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            slide.placeholders[0].text = var_name.replace('value', '').replace('Value', '')

            # Define chart data
            chart_data = CategoryChartData()
            chart_data.categories = conso['cons_year'].index.to_list()
            chart_data.add_series("Consommation Annuelle (kWh)", conso['cons_year'][conso['cons_year'].columns[0]])

            # Define chart type and position
            x, y, cx, cy = Inches(0.5), Inches(2), Inches(9), Inches(5)

            # Add a chart to the slide
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            # Set chart title
            chart.has_title = True
            chart.chart_title.text_frame.text = "Consommation Annuelle (kWh)"

            # Set axes titles
            category_axis = chart.category_axis
            category_axis.has_major_gridlines = True
            category_axis.tick_labels.font.bold = True

            value_axis = chart.value_axis
            value_axis.has_major_gridlines = True
            value_axis.tick_labels.font.bold = True
            value_axis.tick_labels.size = 12

    # Save the presentation
    prs.save(dir + "Graphiques.pptx")

def xl_correction_climatique(conso, meteo, dir):
    '''Utilise la sortie de la fonction consumption_from_index'''
    from openpyxl import load_workbook
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData, XySeriesData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE, XL_LEGEND_POSITION
    from pptx.util import Inches
    import numpy as np
    from scipy import stats

    # INITIALISATION POWERPOINT
    prs = Presentation()

    var_name = conso['cons_day'].columns[0]

    # On calcul la conso par DJU pour le pas de temps de base qui est le JOUR
    meteo.index = meteo.index.strftime('%Y-%m-%d')
    # Concatenation des 2 dataframes en s'assurant de l'egalité de taille
    df_conc = pd.concat([conso['cons_day'], meteo], axis=1)
    df_conc = df_conc.dropna(how='any')

    ## ------------------------------------------------------------------
    # SLIDE SIGNATURE ENERGETIQUE AVEC DJ PAR JOUR
    slide_layout = prs.slide_layouts[5]  # 5 is the index for Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.placeholders[0].text = var_name.replace('value', '').replace('Value', '')

    # Define chart data
    chart_data = XyChartData()
    cd = chart_data.add_series('Mesures', number_format=None)
    x_values_list = df_conc['DJ']
    y_values_list = df_conc[df_conc.columns[0]]
    for x, y in list(zip(x_values_list, y_values_list)):
        cd.add_data_point(x, y, number_format=None)

    # On cree la droite de regression lineaire
    slope, intercept, rvalue, pvalue, stderr = stats.linregress(x_values_list, y_values_list)
    reg_x_values = np.arange(0,max(x_values_list),0.5)
    reg_y_values = slope * reg_x_values + intercept
    reg_series = chart_data.add_series('Régression linéaire')
    for x, y in zip(reg_x_values, reg_y_values):
        reg_series.add_data_point(x, y)

    # Define chart type and position
    x, y, cx, cy = Inches(0.5), Inches(2), Inches(9), Inches(5)

    # Add a chart to the slide
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data
    ).chart

    # Set chart title
    chart.has_title = True
    chart.chart_title.text_frame.text = "Signature Energétique"
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.tick_labels.font.bold = False
    value_axis.tick_labels.size = 10
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    xy_series = chart.series[0]
    xy_series.marker.size = 10

    reg_series = chart.series[1]
    reg_series.marker.style = XL_MARKER_STYLE.CIRCLE
    reg_series.marker.size = 8
    ## FIN SLIDE INITIALE SIGNATURE ENERGETIQUE
    #============================================================================

    ## ------------------------------------------------------------------
    # SLIDE SIGNATURE ENERGETIQUE AVEC RAYONNEMENT SOLAIRE PAR JOUR
    slide_layout = prs.slide_layouts[5]  # 5 is the index for Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.placeholders[0].text = var_name.replace('value', '').replace('Value', '')

    # Define chart data
    chart_data = XyChartData()
    cd = chart_data.add_series('Mesures', number_format=None)
    x_values_list = df_conc['Ray']
    y_values_list = df_conc[df_conc.columns[0]]
    for x, y in list(zip(x_values_list, y_values_list)):
        cd.add_data_point(x, y, number_format=None)

    # Define chart type and position
    x, y, cx, cy = Inches(0.5), Inches(2), Inches(9), Inches(5)

    # Add a chart to the slide
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data
    ).chart

    # Set chart title
    chart.has_title = True
    chart.chart_title.text_frame.text = "Corrélation Rayonnement Solaire ? (kWh global/m2.jour)"
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.tick_labels.font.bold = False
    value_axis.tick_labels.size = 10
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    xy_series = chart.series[0]
    xy_series.marker.size = 10

    ## FIN SLIDE RAYONNEMENT SOLAIRE
    # ============================================================================


    # SUITE ANALYSE
    # SI LE DJ EST NUL, ON MET TOUT A ZERO
    df_conc[df_conc['DJ'] < 0.1] = 0
    df_conc['Corr'] = df_conc[df_conc.columns[0]].div(df_conc[df_conc.columns[1]].replace(0, np.nan))
    df_conc.index = pd.to_datetime(df_conc.index)

    # On check d'abord le range des dates
    total_delta = conso['cons_brut'].index[-1] - conso['cons_brut'].index[0]

    # On initialise les tests
    ae_days, ae_month = False,False

    # Si l'interval est inferieure a un mois, on autorise les graphes data brut et par jour
    if total_delta.days > 90:
        ae_month = True
    if total_delta.days <= 31:
        ae_days = True

    if ae_days:

        df_conc = df_conc.dropna(how='any')

        # SLIDE EVOLUTION CONSOMMATION CORRIGE
        slide_layout = prs.slide_layouts[5]  # 5 is the index for Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        slide.placeholders[0].text = var_name.replace('value', '').replace('Value', '')

        # Define chart data
        chart_data = CategoryChartData()
        chart_data.categories = df_conc.index.to_list()
        chart_data.add_series("Consommation Journalière Corrigée (kWh/DJ)", df_conc['Corr'].to_list())

        # Define chart type and position
        x, y, cx, cy = Inches(0.5), Inches(2), Inches(9), Inches(5)

        # Add a chart to the slide
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart

        # Set chart title
        chart.has_title = True
        chart.chart_title.text_frame.text = "Consommation Journalière Corrigée (kWh/DJ)"

        # Set axes titles
        category_axis = chart.category_axis
        category_axis.has_major_gridlines = True
        category_axis.tick_labels.font.bold = True

        value_axis = chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.tick_labels.font.bold = True
        value_axis.tick_labels.size = 12

    if ae_month:

        df_conc = df_conc.fillna(0)
        df_by_months = df_conc.resample('m').sum()
        df_by_months.index = df_by_months.index.strftime('%B-%Y')

        # Add a slide
        slide_layout = prs.slide_layouts[5]  # 5 is the index for Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        slide.placeholders[0].text = var_name.replace('value','').replace('Value','')

        # Define chart data
        chart_data = CategoryChartData()
        chart_data.categories = df_by_months.index.to_list()
        chart_data.add_series("Consommation Mensuelle Corrigée (kWh/DJ)", df_by_months['Corr'])

        # Define chart type and position
        x, y, cx, cy = Inches(0.5), Inches(2), Inches(9), Inches(5)

        # Add a chart to the slide
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart

        # Set chart title
        chart.has_title = True
        chart.chart_title.text_frame.text = "Consommation Mensuelle Corrigée (kWh/DJ)"

        # Set axes titles
        category_axis = chart.category_axis
        category_axis.has_major_gridlines = True
        category_axis.tick_labels.font.bold = False

        value_axis = chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.tick_labels.font.bold = False
        value_axis.tick_labels.size = 12

    # Save the presentation
    prs.save(dir + "Graphiques_Correction_Climatique.pptx")

def write_conso_files(nom, conso, dir):

    # On ecrit les consos dans les fichiers csv
    # On verifie d'abord l'existence du repertoire
    if not os.path.isdir(dir + '\\CONSOMMATIONS\\'):
        os.mkdir(dir + '\\CONSOMMATIONS\\')

    conso['cons_brut'].to_csv(dir + '\\CONSOMMATIONS\\' + 'EXTRACT_BRUT_' + nom + '.csv', index=True)
    conso['cons_day'].to_csv(dir + '\\CONSOMMATIONS\\' + 'EXTRACT_JOURS_' + nom + '.csv', index=True)
    conso['cons_month'].to_csv(dir + '\\CONSOMMATIONS\\' + 'EXTRACT_MOIS_' + nom + '.csv', index=True)
    conso['cons_year'].to_csv(dir + '\\CONSOMMATIONS\\' + 'EXTRACT_ANNEE_' + nom + '.csv', index=True)

    # Un excel avec tous les types de conso egalement
    writer = pd.ExcelWriter(dir + '\\CONSOMMATIONS\\' + '\\EXTRACT_CONSOS_' + nom + '.xlsx', engine='xlsxwriter',
                            engine_kwargs={'options': {'strings_to_numbers': True}})
    conso['cons_brut'].to_excel(writer, index=True, sheet_name="Brut")
    conso['cons_day'].to_excel(writer, index=True, sheet_name="Jours")
    conso['cons_month'].to_excel(writer, index=True, sheet_name="Mois")
    conso['cons_year'].to_excel(writer, index=True, sheet_name="Annees")

    writer.close()

def make_slide_plot(liste_df, prs):
    """return a chart object"""
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE, XL_LEGEND_POSITION
    from pptx.util import Inches

    for df in liste_df:

        # ON commence par retirer les NaN
        df = df.dropna(how='any')
        var_name = df.columns[0]

        # Add a slide
        slide_layout = prs.slide_layouts[5]  # 5 is the index for Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        slide.placeholders[0].text = var_name.replace('value', '').replace('Value', '')
    
        # Define chart data
        chart_data = CategoryChartData()
        chart_data.categories(df.index.to_list())
        chart_data.add_series(var_name,df[df.columns[0]])

        # Define chart type and position
        x, y, cx, cy = Inches(0.5), Inches(2), Inches(9), Inches(5)

        # Add a chart to the slide
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
        ).chart
    
        # Set chart title
        chart.has_title = True
        chart.chart_title.text_frame.text = var_name
    
        # Set axes titles
        category_axis = chart.category_axis
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.bold = False
    
        value_axis = chart.value_axis
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.bold = False
        value_axis.tick_labels.size = 12

    return prs

def main():

    dir_data = 'C:\\EXTRACTIONS\\JTI_DATA\\'
    dir_write = 'C:\\NETTOYEUR\\'
    # Creation du dossier de travail
    if not os.path.isdir(dir_write):
        os.mkdir(dir_write)

    df_brut = pd.read_csv(dir_data+'Test_JTI.csv')
    df_brut = df_brut.set_index('ts')
    df_brut.index = pd.to_datetime(df_brut.index)

    df_brut = df_brut.drop(['quality'], axis=1, errors='ignore')

    df_out, df_in = detect_outliers_iqr(df_brut,'tagValue')

    conso = consumption_from_index(df_in)

    #xl_consumption(conso,dir_write)


    # Ecriture de la figure de synthese
    fig, (ax1, ax2) = plt.subplots(2,1)
    ax1.scatter(df_in.index,df_in['tagValue'], s=50, marker='o',color='black')
    ax1.scatter(df_out.index, df_out['tagValue'], s=50, marker='o', color='red',label='discarded data')
    ax2.scatter(df_in.index, df_in['tagValue'], s=50, marker='o', color='green',label='after cleaning')
    ax1.legend(fontsize=12)
    ax2.legend(fontsize=12)
    #plt.show()
    plt.suptitle("tagName : " + 'tagValue')
    plt.savefig(dir_write+'\\tagValue.png')
    plt.close()
    create_pdf(dir_write + '\\Rapport_Nettoyage_' + 'tagValue'+'.pdf', dir_write+'\\tagValue.png')


    print('fin du main')


if __name__ == '__main__':
    main()
